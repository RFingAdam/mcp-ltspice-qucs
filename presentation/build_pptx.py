"""Generate the investor pitch deck as OVERVIEW.pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent
ASSETS = ROOT / "assets"

NAVY = RGBColor(0x1F, 0x2A, 0x44)
TEAL = RGBColor(0x1A, 0x8A, 0x8A)
ORANGE = RGBColor(0xE7, 0x75, 0x4A)
GRAY = RGBColor(0x6B, 0x72, 0x80)


def _set_text(tf, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_bullets(tf, bullets, size=14, color=NAVY):
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(8)


def _add_title(slide, title, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6))
    _set_text(tb.text_frame, title, size=28, bold=True, color=NAVY)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.33), Inches(0.4))
        _set_text(sb.text_frame, subtitle, size=14, color=GRAY)


def slide_cover(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(1.5))
    _set_text(
        tb.text_frame, "mcp-ltspice-qucs", size=60, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)
    )

    sb = s.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(12), Inches(1.0))
    _set_text(
        sb.text_frame,
        "An agentic MCP suite for RF, analog, and SMPS-EMC design",
        size=22,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )

    tag = s.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(12), Inches(0.6))
    _set_text(tag.text_frame, "Spec → compliant deliverable in seconds.", size=18, color=TEAL)

    foot = s.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(12), Inches(0.4))
    _set_text(
        foot.text_frame,
        "v0.2.0  •  427 tests passing  •  Investor briefing",
        size=12,
        color=RGBColor(0xCB, 0xD5, 0xE1),
    )


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "The problem", "Hardware design today is slow, expensive, and reactive")

    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.0))
    _add_bullets(
        tb.text_frame,
        [
            "3–10 days per filter design — spreadsheet math → ADS / Cadence "
            "($30k–$200k per seat) → simulator → respin",
            "60–80 % of products fail CISPR / FCC pre-compliance on first lab visit "
            "— each respin costs $20k–$80k NRE + 2–6 weeks slip",
            "Vendor part substitution is manual: 100-tab Excel files. "
            "Parasitics (SRF, Q, ESR) routinely ignored — silently degrade the design",
            "EMC compliance is reactive, not predicted — the engineer doesn't "
            "know they failed until the lab tells them",
        ],
        size=18,
    )

    tb2 = s.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(12.33), Inches(1.2))
    _set_text(
        tb2.text_frame,
        "Today's flow: engineer → math → simulator → respin.\n"
        "We collapse the loop so an LLM agent does engineer-quality work end-to-end.",
        size=16,
        bold=True,
        color=ORANGE,
    )


def slide_what_we_built(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(
        s, "What we built", "112 MCP tools across four packages, callable by any agent stack"
    )

    s.shapes.add_picture(str(ASSETS / "hero.png"), Inches(0.5), Inches(1.4), width=Inches(12.33))

    cap = s.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.33), Inches(0.4))
    _set_text(
        cap.text_frame,
        "Three categories · 56 flat tools + 56 namespaced aliases · "
        "all backed by peer-reviewed math (Pozar, Erickson, Sedra-Smith, Mancini)",
        size=11,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )


def slide_demo1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(
        s,
        "Demo 1 — RF Low-Pass Filter (1 GHz Butterworth)",
        "Synthesise → analyse → Monte Carlo → LTspice/Touchstone deliverable in 9 seconds",
    )

    s.shapes.add_picture(
        str(ASSETS / "demo1_lpf.png"), Inches(1.67), Inches(1.5), width=Inches(10.0)
    )

    tb = s.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(12.33), Inches(2.0))
    _add_bullets(
        tb.text_frame,
        [
            "Spec: passband to 600 MHz at 0.5 dB IL; 30/45/60 dB at 2·/3·/5·fc",
            "Monte Carlo: 65 % yield on ideal ±5 % parts → 99 % with real "
            "Coilcraft + Murata vendor substitution (E24 snap, SRF-aware)",
            "Tangible output: .asc, .s2p, schematic PNG, PDF report — engineer-ready",
        ],
        size=14,
    )


def slide_demo2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(
        s,
        "Demo 2 — SMPS Conducted Emissions vs CISPR 32 Class B",
        "From 'I have a switching converter' to 'here's how close to compliance you are'",
    )

    s.shapes.add_picture(
        str(ASSETS / "demo2_emc.png"), Inches(1.67), Inches(1.4), width=Inches(10.0)
    )

    tb = s.shapes.add_textbox(Inches(0.5), Inches(5.55), Inches(12.33), Inches(1.9))
    _add_bullets(
        tb.text_frame,
        [
            "Modelled the switch-node as a trapezoidal waveform, predicted "
            "all 80 harmonics across the 150 kHz – 30 MHz CISPR conducted band",
            "Designed a 70-dB-at-fsw LC input filter; verified Middlebrook "
            "stability against the converter's input impedance",
            "Worst margin: −2.8 dB — engineer immediately knows: ~3 dB more "
            "attenuation needed (taller ferrite, larger X-cap, or CM choke). "
            "No spin-the-board.",
        ],
        size=14,
    )


def slide_demo3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(
        s,
        "Demo 3 — Active Filter (Sallen-Key, 22 kHz, 4th order)",
        "Anti-aliasing for a 96 kSPS audio ADC, complete with op-amp recommendation",
    )

    s.shapes.add_picture(
        str(ASSETS / "demo3_sallen_key.png"), Inches(1.67), Inches(1.4), width=Inches(10.0)
    )

    tb = s.shapes.add_textbox(Inches(0.5), Inches(5.55), Inches(12.33), Inches(1.9))
    _add_bullets(
        tb.text_frame,
        [
            "Decomposed into two Sallen-Key biquads with Mancini-tabulated "
            "stage Qs (0.541 / 1.307); sized R/C for 1 nF capacitors",
            "Required op-amp GBW: 2.87 MHz. Bundled catalogue auto-screened "
            "8 candidates against GBW, noise (≤ 10 nV/√Hz), offset (≤ 1 mV)",
            "Audio-band rank (cap GBW, sort by noise): ADI LT6275 — "
            "45 MHz · 1.9 nV/√Hz · 130 µV. Per-stage schematics + PDF report rendered.",
        ],
        size=14,
    )


def slide_moat(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "What's hard to replicate", "Data, math, and agent-first design")

    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.5))
    _add_bullets(
        tb.text_frame,
        [
            "Real vendor data baked in — Coilcraft, Johanson, TDK, Murata "
            "L/Q/SRF tables with parasitic-aware substitution. Not catalog scraping.",
            "Peer-reviewed math — Pozar §8.5 RF transforms, Hammerstad-Jensen "
            "microstrip, Erickson §10.4 EMI, Sedra-Smith Sallen-Key, Mancini cascade tables",
            "427 unit tests guard the math, 4 simulator-gated skips, 0 failures",
            "Designed agent-first — every tool returns Envelope[T] with status, "
            "data, warnings, ms-timer. Exactly what an LLM needs to chain reliably.",
            "Composable with the rest of the engineering MCP ecosystem — "
            "antennas → nec2-antenna / openems, PCB layout → pcb-emcopilot, "
            "regulatory → emc-regulations. We sit cleanly in the middle.",
            "v0.2.0 shipped, tagged, GitHub Release public. Reproducible from "
            "git clone + uv run pytest.",
        ],
        size=15,
    )


def slide_why_now(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "Why now")

    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.0))
    _add_bullets(
        tb.text_frame,
        [
            "MCP standardised in late 2024 — build once, every agent stack consumes it",
            "LLM agents have crossed the 'useful for engineering' threshold — "
            "Claude 4.x and GPT-5 reliably orchestrate multi-step tool use",
            "Hardware-design tooling moat is the data + the math, not the wrapper code "
            "— and that's where we've invested",
            "EDA incumbents (Cadence, Synopsys, Keysight) are pricing for the "
            "Fortune-500 IDM market — the long tail of IoT / startup / "
            "consumer-electronics teams is wide open",
        ],
        size=18,
    )

    box = s.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.8))
    _set_text(
        box.text_frame,
        "The hardware team that ships first wins. We make every team's HW lead 5–10× faster.",
        size=18,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )


def slide_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(s, "Roadmap")

    rows = [
        (
            "Now (v0.2.0)",
            "112 tools shipped · RF LPF/HPF/BPF/BSF · SMPS EMC · "
            "active filters · vendor catalogues · LTspice/Touchstone deliverables",
        ),
        (
            "Phase 6",
            "Xyce harmonic-balance · Qucs-S noise extraction · "
            "distributed filter synthesis (hairpin, interdigital, combline) · "
            ".sch ↔ .asc conversion bridge",
        ),
        (
            "Phase 7",
            "Real-simulator integration tests · async runner with heartbeat · "
            "Sobol-index sensitivity · correlated-tolerance Monte Carlo",
        ),
        (
            "Beyond",
            "Persistent ngspice (10–50× MC speedup) · FilterDesign threading "
            "through pipeline · live vendor-fetch agents (Coilcraft / Murata APIs)",
        ),
    ]
    y = Inches(1.6)
    for label, body in rows:
        lbl = s.shapes.add_textbox(Inches(0.5), y, Inches(2.3), Inches(0.8))
        _set_text(lbl.text_frame, label, size=15, bold=True, color=TEAL)
        body_box = s.shapes.add_textbox(Inches(2.9), y, Inches(10.0), Inches(1.2))
        _set_text(body_box.text_frame, body, size=13, color=NAVY)
        y += Inches(1.05)


def slide_close(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    title = s.shapes.add_textbox(Inches(0.7), Inches(2.6), Inches(12), Inches(1.0))
    _set_text(
        title.text_frame,
        "Validated. Shipped. Composable.",
        size=44,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
    )

    sub = s.shapes.add_textbox(Inches(0.7), Inches(3.8), Inches(12), Inches(0.7))
    _set_text(
        sub.text_frame,
        "github.com/RFingAdam/mcp-ltspice-qucs   ·   v0.2.0",
        size=18,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )

    foot = s.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(12), Inches(0.5))
    _set_text(
        foot.text_frame,
        "Ready to demo live.",
        size=16,
        color=RGBColor(0xCB, 0xD5, 0xE1),
        align=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)
    slide_problem(prs)
    slide_what_we_built(prs)
    slide_demo1(prs)
    slide_demo2(prs)
    slide_demo3(prs)
    slide_moat(prs)
    slide_why_now(prs)
    slide_roadmap(prs)
    slide_close(prs)
    out = ROOT / "OVERVIEW.pptx"
    prs.save(out)
    print(f"Wrote {out} ({out.stat().st_size // 1024} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
